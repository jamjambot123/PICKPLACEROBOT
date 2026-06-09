# -*- coding: utf-8 -*-
"""
최종 포스터 PPT 생성 스크립트
템플릿의 기존 레이아웃/디자인을 보존하면서 내용을 채워 넣습니다.
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import copy
from lxml import etree

# ── 경로 설정 ──────────────────────────────────────────────────
TEMPLATE = r'd:\dongwon\PICKPLACEROBOT\[4학년] (양식)2026년 산학프로젝트 포스터_260512.ver.pptx'
OUTPUT   = r'd:\dongwon\PICKPLACEROBOT\[완성본] 2026년 산학프로젝트 포스터.pptx'

prs = Presentation(TEMPLATE)
slide = prs.slides[0]

# ── 도우미 함수 ────────────────────────────────────────────────
def find_shape(name):
    """이름으로 shape 찾기"""
    for s in slide.shapes:
        if s.name == name:
            return s
    raise ValueError(f"Shape '{name}' not found")


def clear_and_set_text(shape, text):
    """Shape의 텍스트를 단일 문단으로 교체 (첫 번째 run의 서식 보존)"""
    tf = shape.text_frame
    # 첫 번째 paragraph의 첫 번째 run에서 서식 복사
    first_para = tf.paragraphs[0]
    fmt_xml = None
    if first_para.runs:
        fmt_xml = copy.deepcopy(first_para.runs[0]._r.find(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr'))

    # 모든 기존 paragraph 제거 후 새 텍스트 설정
    p_elements = tf._txBody.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}p')
    for p in p_elements[1:]:
        tf._txBody.remove(p)

    # 첫 번째 paragraph 정리 후 텍스트 설정
    first_p = p_elements[0]
    for r in first_p.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}r'):
        first_p.remove(r)

    nsmap = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    new_r = etree.SubElement(first_p, f'{nsmap}r')
    if fmt_xml is not None:
        new_r.insert(0, fmt_xml)
    new_t = etree.SubElement(new_r, f'{nsmap}t')
    new_t.text = text


def set_multiline(shape, lines, section_font_size=None, body_font_size=None,
                  section_bold=True, alignment=None):
    """
    Shape에 여러 줄의 텍스트를 설정합니다.
    [섹션 제목]은 볼드, 나머지는 일반체로 설정합니다.
    """
    tf = shape.text_frame
    nsmap = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

    # 기존 서식 참조 저장
    first_para = tf.paragraphs[0]
    orig_pPr = first_para._p.find(f'{nsmap}pPr')
    orig_rPr = None
    if first_para.runs:
        orig_rPr = first_para.runs[0]._r.find(f'{nsmap}rPr')

    # 모든 기존 paragraph 제거
    p_elements = tf._txBody.findall(f'{nsmap}p')
    for p in p_elements:
        tf._txBody.remove(p)

    for i, line in enumerate(lines):
        new_p = etree.SubElement(tf._txBody, f'{nsmap}p')

        # paragraph 속성 복사
        if orig_pPr is not None:
            new_pPr = copy.deepcopy(orig_pPr)
            new_p.insert(0, new_pPr)

        if alignment is not None:
            pPr = new_p.find(f'{nsmap}pPr')
            if pPr is None:
                pPr = etree.SubElement(new_p, f'{nsmap}pPr')
                new_p.insert(0, pPr)
            pPr.set('algn', alignment)

        # 섹션 제목인지 판별
        is_section = line.startswith('[') or line.startswith('■')
        is_empty = len(line.strip()) == 0

        new_r = etree.SubElement(new_p, f'{nsmap}r')

        # run 속성 설정
        new_rPr = etree.SubElement(new_r, f'{nsmap}rPr')
        new_rPr.set('lang', 'ko-KR')
        new_rPr.set('altLang', 'en-US')

        if is_section and section_bold:
            new_rPr.set('b', '1')
        
        # 폰트 크기 설정
        if is_section and section_font_size:
            new_rPr.set('sz', str(section_font_size))
        elif body_font_size and not is_empty:
            new_rPr.set('sz', str(body_font_size))

        new_t = etree.SubElement(new_r, f'{nsmap}t')
        new_t.text = line


def add_table_to_slide(slide, rows_data, left, top, width, height):
    """슬라이드에 테이블을 추가합니다."""
    num_rows = len(rows_data)
    num_cols = len(rows_data[0])
    
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # 열 너비 설정
    col_widths = [Emu(3200000), Emu(3600000), Emu(3600000), Emu(3600000)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for r_idx, row_data in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = cell_text
            
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(18)
                    run.font.name = '맑은 고딕'
                    if r_idx == 0:  # 헤더 행
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    elif c_idx == 0:  # 첫 번째 열
                        run.font.bold = True

            # 헤더 행 배경색
            if r_idx == 0:
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = RGBColor(0x00, 0x32, 0x6E)  # 진한 남색
            # 데이터 행 번갈아 배경색
            elif r_idx % 2 == 1:
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF7)

    return table_shape


# ═══════════════════════════════════════════════════════════════
# 1. 제목 (연구주제) - TextBox 21
# ═══════════════════════════════════════════════════════════════
title_shape = find_shape('TextBox 21')
clear_and_set_text(title_shape,
    'AI 기반 비대칭 S-Curve 궤적 최적화 및 인풋쉐이핑을 활용한\n'
    '반도체 픽앤플레이스 갠트리 진동 억제 시스템')

# ═══════════════════════════════════════════════════════════════
# 2. 팀 정보 수정
# ═══════════════════════════════════════════════════════════════
# 학과 - TextBox 16
dept_shape = find_shape('TextBox 16')
clear_and_set_text(dept_shape, '전자공학과')

# 지도교수 - TextBox 17
prof_shape = find_shape('TextBox 17')
clear_and_set_text(prof_shape, '송동섭 교수님')

# 산학기관 - TextBox 18
company_shape = find_shape('TextBox 18')
clear_and_set_text(company_shape, '호서반도체')

# 산학기관2 - TextBox 19
company2_shape = find_shape('TextBox 19')
clear_and_set_text(company2_shape, '삼성전자')

# 팀원 - TextBox 20
members_shape = find_shape('TextBox 20')
tf = members_shape.text_frame
nsmap = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
# 기존 내용 유지 (이미 입력됨)

# ═══════════════════════════════════════════════════════════════
# 3. 주제 소개 영역 (사각형: 둥근 모서리 8 - 프로젝트 제목)
# ═══════════════════════════════════════════════════════════════
subtitle_shape = find_shape('사각형: 둥근 모서리 8')
clear_and_set_text(subtitle_shape,
    'AI 기반 비대칭 S-Curve 궤적 최적화 및 인풋쉐이핑을 활용한 반도체 픽앤플레이스 갠트리 진동 억제 시스템')

# ═══════════════════════════════════════════════════════════════
# 4. 프로젝트 제작동기 영역 (TextBox 13)
# ═══════════════════════════════════════════════════════════════
motivation_shape = find_shape('TextBox 13')
clear_and_set_text(motivation_shape, '프로젝트 제작동기 및 연구 배경')

# ═══════════════════════════════════════════════════════════════
# 5. 본문 상단 영역 (사각형: 둥근 모서리 10) - 연구배경 + 시스템모델 + 실험결과
# ═══════════════════════════════════════════════════════════════
body_shape = find_shape('사각형: 둥근 모서리 10')

body_lines = [
    '■ 1. 연구 배경',
    '',
    '• 반도체 Pick & Place 공정에서 갠트리 고속 이동 시 잔류 진동 발생',
    '• 진동이 정착될 때까지 대기 → UPH(생산성) 저하의 직접적 원인',
    '• 해결: 모션 프로파일 최적화 + 인풋쉐이핑(ZVD)으로 진동 원천 억제',
    '',
    '',
    '■ 2. 시스템 모델',
    '',
    '• 2-Mass Model: Motor(Perfect Servo) → Spring-Damper → Load',
    '• RK4 수치적분, dt=0.1ms, fn=10Hz, ζ=0.05',
    '• 디지털 트윈: Three.js 3D + Chart.js 실시간 시각화',
    '',
    '',
    '■ 3. 실험 결과 (정량적 비교)',
    '',
    '                                    Trapezoidal        AS-Curve         AS-Curve+ZVD',
    '   잔류 진동                       17.47mm            24.66mm           4.29mm',
    '   정착 시간                       2.660s              2.815s             0.197s',
    '   UPH                               676                  639               9,137',
]

set_multiline(body_shape, body_lines, 
              section_font_size=1800, body_font_size=1600,
              section_bold=True)

# ═══════════════════════════════════════════════════════════════
# 6. 프로젝트 진행상황 (TextBox 14) 
# ═══════════════════════════════════════════════════════════════
progress_shape = find_shape('TextBox 14')
clear_and_set_text(progress_shape, '시스템 모델 및 실험 결과 상세')

# ═══════════════════════════════════════════════════════════════
# 7. 하단 제목 (사각형: 둥근 모서리 11) - "작품 구현 및 결과"
# ═══════════════════════════════════════════════════════════════
result_title_shape = find_shape('사각형: 둥근 모서리 11')
clear_and_set_text(result_title_shape, '핵심 성과 및 결론')

# ═══════════════════════════════════════════════════════════════
# 8. 하단 설명 (TextBox 15) - 작품 구현 및 결과 설명 레이블
# ═══════════════════════════════════════════════════════════════
result_desc_shape = find_shape('TextBox 15')
clear_and_set_text(result_desc_shape, '핵심 성과 및 향후 과제')

# ═══════════════════════════════════════════════════════════════
# 9. 실험 결과 테이블 추가
# ═══════════════════════════════════════════════════════════════
table_data = [
    ['구분',         'Trapezoidal', 'AS-Curve', 'AS-Curve+ZVD'],
    ['잔류 진동',    '17.47mm',     '24.66mm',  '4.29mm'],
    ['정착 시간',    '2.660s',      '2.815s',   '0.197s'],
    ['UPH',          '676',         '639',      '9,137'],
]

# 테이블 위치: 본문 영역 아래쪽 (사각형: 둥근 모서리 10 내부 하단)
table_left = Emu(1500000)
table_top  = Emu(23500000)  # 본문 영역 아래
table_width = Emu(14000000)
table_height = Emu(3200000)

add_table_to_slide(slide, table_data, table_left, table_top, table_width, table_height)

# ═══════════════════════════════════════════════════════════════
# 10. 핵심 성과 + 결론 텍스트 박스 추가 (하단 영역)
# ═══════════════════════════════════════════════════════════════
from pptx.util import Inches

# 핵심 성과 텍스트 박스
result_left = Emu(500000)
result_top = Emu(33200000)
result_width = Emu(17200000)
result_height = Emu(8000000)

txBox = slide.shapes.add_textbox(result_left, result_top, result_width, result_height)
tf = txBox.text_frame
tf.word_wrap = True

nsmap = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

result_lines = [
    ('■ 4. 핵심 성과', True, 1800),
    ('', False, 1400),
    ('• 잔류 진동 75.4% 감소 (17.47mm → 4.29mm)', False, 1600),
    ('• 정착 시간 92.6% 단축 (2.660s → 0.197s)', False, 1600),
    ('• UPH 1,252% 향상 (676 → 9,137)', False, 1600),
    ('', False, 1400),
    ('', False, 1400),
    ('■ 5. 결론 및 향후 과제', True, 1800),
    ('', False, 1400),
    ('• ZVD 인풋쉐이핑은 약 100ms의 지연으로 2.5초의 진동 대기를 제거', False, 1600),
    ('• 비대칭 S-Curve + ZVD 조합으로 최적의 진동 억제 성능 달성', False, 1600),
    ('• 향후 실물 하드웨어(STM32 + IMU) 검증 필요', False, 1600),
    ('• PSO 알고리즘을 통한 파라미터 자동 최적화 고도화 계획', False, 1600),
]

for i, (text, bold, font_size) in enumerate(result_lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size / 100)
    run.font.name = '맑은 고딕'
    run.font.bold = bold
    if bold:
        run.font.color.rgb = RGBColor(0x00, 0x32, 0x6E)


# ═══════════════════════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f'✅ 포스터 생성 완료: {OUTPUT}')
