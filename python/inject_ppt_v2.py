import os
from pptx import Presentation
from pptx.util import Pt

def inject_poster():
    template_path = '[4학년] (양식)2026년 산학프로젝트 포스터_260512.ver.pptx'
    if not os.path.exists(template_path):
        return
        
    prs = Presentation(template_path)
    slide = prs.slides[0]
    
    title_text = "인풋쉐이핑 및 비대칭 S-curve 기반 픽앤플레이스 갠트리 진동 제어 시스템 및 디지털 트윈"
    content_text = """[1. 연구 배경 및 목적]
- 반도체 픽앤플레이스 장비의 UPH(생산성) 향상을 위해 고속 이송 시 로봇팔 끝단에 잔류진동 발생, 플레이스 정확도 저하.
- 본 연구는 STM32 기반 1축 갠트리와 로봇팔 시스템에 인풋쉐이핑 및 비대칭 S-curve를 적용하여 정착시간을 단축하고 UPH 향상을 목표로 함.

[2. 시스템 구성 및 제어 알고리즘]
- 하드웨어: STM32(X축 갠트리 모션 제어), Arduino(다관절 로봇팔 제어), MPU6050 IMU(끝단 잔류진동 측정)
- 모션 제어: 
  ① Trapezoidal (기준) 
  ② 비대칭 S-Curve (저크 저감) 
  ③ 인풋쉐이핑 + 비대칭 S-Curve (고유진동수 기반 잔류진동 원천 억제)

[3. 디지털 트윈 및 검증 결과]
- 3D 웹 디지털 트윈을 구축하여 파라미터(V_max, A_max, Beta 등) AI 오토튜닝 구현.
- IMU 센서 실측 결과, 제어 미적용(Trapezoidal) 대비 제어 적용(Shaped) 시 잔류진동(g) 수렴 시간이 대폭 단축되어 실질적 UPH 향상을 입증함.
"""
    
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    if text_shapes:
        text_shapes[0].text_frame.text = title_text
        if len(text_shapes) > 1:
            text_shapes[1].text_frame.text = content_text
        else:
            txBox = slide.shapes.add_textbox(Pt(50), Pt(150), prs.slide_width - Pt(100), prs.slide_height - Pt(200))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.text = content_text
            for p in tf.paragraphs: p.font.size = Pt(28)
            
    prs.save(template_path)


def inject_presentation():
    template_path = '[4학년] (양식)2026년 산학프로젝트 최종보고 발표자료.pptx'
    if not os.path.exists(template_path):
        return
        
    prs = Presentation(template_path)
    
    slide_contents = [
        # Slide 0: Title
        ("인풋쉐이핑 및 비대칭 S-curve 기반\n픽앤플레이스 갠트리 진동 제어 시스템 및 디지털 트윈", "2026년 산학프로젝트 최종보고"),
        # Slide 1: 목차 
        ("목차", "1. 연구 배경 및 목적\n2. 시스템 하드웨어 구성\n3. 핵심 제어 알고리즘 (AS-Curve & Input Shaping)\n4. 디지털 트윈 시뮬레이션\n5. 실험 결과 및 결론 (UPH 향상)"),
        # Slide 2: Background
        ("1. 연구 배경: 반도체 장비의 UPH 딜레마", "• 반도체 패키징 공정에서 생산성(UPH) 향상을 위해 고속 이송 필수\n• 고속 이동 시 갠트리 탑재 로봇팔 끝단에 심각한 '잔류 진동' 발생\n• 진동 수렴을 기다리는 '정착 시간'으로 인해 실질적 UPH 저하 및 플레이스 오차 발생"),
        # Slide 3: Hardware
        ("2. 시스템 하드웨어 및 통신 구조", "• 제어부:\n  - STM32: 갠트리 X축 모션 제어 및 프로파일 생성\n  - Arduino: 다관절 로봇팔 제어\n  - UART 시리얼 통신으로 픽앤플레이스 동작 동기화\n• 센서부:\n  - MPU6050 IMU: 로봇팔 끝단에 부착하여 잔류 진동 가속도(g) 실시간 측정"),
        # Slide 4: Algorithms
        ("3. 핵심 제어 알고리즘: 단계별 모션 프로파일", "① Trapezoidal: 일정 가감속, 높은 저크 발생 (비교 기준)\n② 비대칭 S-Curve: 가속 구간은 빠르게, 감속 구간은 부드럽게 설정하여 물리적 충격 저감\n③ Input Shaping + AS-Curve: 로봇의 고유진동수를 기반으로 ZV/ZVD 쉐이퍼를 설계하여 진동 자체를 원천 상쇄 (Cancel-out)"),
        # Slide 5: Digital Twin
        ("4. 디지털 트윈 시뮬레이터 구축", "• 이론 및 하드웨어 실험을 넘어 웹 기반 3D 디지털 트윈 구축\n• 고해상도 실시간 분석 차트 (위치, 속도, 가속도, 진동 오차)\n• AI 오토튜닝 (PSO 알고리즘) 적용:\n  - 갠트리의 한계 속도 내에서 최적의 파라미터(V_max, A_max, Beta, Gamma)를 자동 탐색"),
        # Slide 6: Results
        ("5. 검증 결과 및 결론", "• [IMU 실측 결과] 제어 미적용(Trapezoidal) 대비 제어 적용(Shaped) 시 잔류진동 정착시간 획기적 단축\n• [효과] 정착시간 단축분을 UPH 향상률로 환산 시 O% 향상 달성\n• [결론] 오프라인 디지털 트윈 튜닝과 실제 하드웨어 제어의 결합으로 반도체 장비 생산성 극대화 솔루션 제시")
    ]
    
    for i, slide in enumerate(prs.slides):
        if i < len(slide_contents):
            title, content = slide_contents[i]
            
            text_shapes = [s for s in slide.shapes if s.has_text_frame]
            title_shape = None
            body_shape = None
            
            for shape in slide.placeholders:
                if shape.placeholder_format.type in [1, 3]: title_shape = shape
                elif shape.placeholder_format.type in [2, 7]: body_shape = shape
                    
            if title_shape and title_shape.has_text_frame:
                title_shape.text_frame.text = title
            elif len(text_shapes) > 0:
                text_shapes[0].text_frame.text = title
                
            if body_shape and body_shape.has_text_frame:
                body_shape.text_frame.text = content
            elif len(text_shapes) > 1:
                text_shapes[1].text_frame.text = content
                
            if not title_shape and not body_shape and len(text_shapes) == 0:
                txBox = slide.shapes.add_textbox(Pt(50), Pt(50), prs.slide_width - Pt(100), Pt(50))
                txBox.text_frame.text = title
                txBox2 = slide.shapes.add_textbox(Pt(50), Pt(120), prs.slide_width - Pt(100), prs.slide_height - Pt(150))
                txBox2.text_frame.text = content

    prs.save(template_path)

if __name__ == '__main__':
    inject_poster()
    inject_presentation()
    print("Project proposal content injected.")
