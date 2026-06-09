import os
from pptx import Presentation
from pptx.util import Pt

def inject_poster():
    template_path = '[4학년] (양식)2026년 산학프로젝트 포스터_260512.ver.pptx'
    if not os.path.exists(template_path):
        return
        
    prs = Presentation(template_path)
    slide = prs.slides[0]
    
    title_text = "인풋쉐이핑 및 비대칭 S-curve 기반 픽앤플레이스 진동 제어 AI 시뮬레이터"
    content_text = """[1. 연구 배경 및 목적]
- 반도체 픽앤플레이스 장비의 생산성(UPH) 향상을 위해 고속 이송 시 발생하는 잔류진동 억제 필수.
- 하드웨어 제작 전, 물리 동역학 방정식 기반의 '디지털 트윈 시뮬레이터'를 개발하여 최적의 모션 제어 알고리즘을 설계하고 검증함.

[2. 핵심 제어 알고리즘]
- ① Trapezoidal (기준 프로파일)
- ② 비대칭 S-Curve: 가속/감속 구간의 물리적 충격(Jerk) 최소화
- ③ Input Shaping: 파동 상쇄 간섭 원리를 이용해 잔류진동 원천 억제 (ZV/ZVD 쉐이퍼)

[3. 시뮬레이션 및 AI 최적화 (PSO)]
- Three.js 기반 3D 웹 디지털 트윈을 구축하여 궤적별 실시간 진동 오차 시각화.
- 입자군집최적화(PSO) 알고리즘을 백그라운드에 적용하여, 진동이 0이 되면서 가장 빠른 정착시간을 내는 파라미터(V_max, A_max, Beta 등) 자동 탐색.

[4. 연구 결과 및 기대 효과]
- 시뮬레이션 결과, 제어 미적용(Trapezoidal) 대비 제어 적용(Shaped) 시 잔류진동 수렴 시간이 획기적으로 단축됨.
- 오프라인 시뮬레이션을 통한 사전 AI 튜닝으로 반도체 장비 셋업 시간 단축 및 UPH 극대화 기여."""
    
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    if text_shapes:
        text_shapes[0].text_frame.text = title_text
        if len(text_shapes) > 1:
            text_shapes[1].text_frame.text = content_text
            for p in text_shapes[1].text_frame.paragraphs:
                p.font.size = Pt(18)  # 포스터 글자 크기 축소
        else:
            txBox = slide.shapes.add_textbox(Pt(50), Pt(150), prs.slide_width - Pt(100), prs.slide_height - Pt(200))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.text = content_text
            for p in tf.paragraphs: p.font.size = Pt(18)
            
    prs.save(template_path)


def inject_presentation():
    template_path = '[4학년] (양식)2026년 산학프로젝트 최종보고 발표자료.pptx'
    if not os.path.exists(template_path):
        return
        
    prs = Presentation(template_path)
    
    slide_contents = [
        # Slide 0: Title
        ("인풋쉐이핑 및 비대칭 S-curve 기반\n픽앤플레이스 진동 제어 AI 디지털 트윈 시뮬레이터", "2026년 산학프로젝트 최종보고"),
        # Slide 1: 목차 
        ("목차", "1. 연구 배경 및 목적\n2. 동역학 모델링 및 핵심 제어 알고리즘\n3. 디지털 트윈 시뮬레이터 개발\n4. AI 파라미터 자동 최적화 (PSO)\n5. 시뮬레이션 결과 및 결론"),
        # Slide 2: Background
        ("1. 연구 배경: 반도체 장비의 UPH 딜레마", "• 생산성(UPH) 향상을 위해 고속 이송 필수\n• 고속 이동 시 로봇팔 끝단에 심각한 '잔류 진동(Residual Vibration)' 발생\n• 진동 수렴을 기다리는 '정착 시간(Settling Time)' 증가로 실질적 UPH 저하\n• 하드웨어 셋업 전, 이를 예측하고 최적화할 오프라인 시뮬레이터 부재"),
        # Slide 3: Algorithm
        ("2. 핵심 제어 알고리즘: 단계별 모션 프로파일", "① Trapezoidal (사다리꼴): 일정 가감속, 높은 저크(Jerk) 발생 (비교 기준)\n② 비대칭 S-Curve: 가속 구간은 빠르게, 감속 구간은 부드럽게 설정하여 물리적 충격 저감\n③ Input Shaping: 로봇의 고유진동수를 기반으로 ZV 쉐이퍼를 설계하여 파동의 상쇄 간섭을 통해 진동 원천 제거"),
        # Slide 4: Simulation
        ("3. 3D 디지털 트윈 시뮬레이터 개발", "• 물리 동역학 방정식(ODE)을 웹(Web) 환경의 3D 엔진(Three.js)에 이식\n• 픽앤플레이스 로봇의 움직임과 진동을 실시간 교차 검증\n• 4대 핵심 분석 차트 탑재:\n  - 위치 추적, 목표 속도 궤적, 가속도 및 저크, 잔류 진동 오차 시각화"),
        # Slide 5: AI Optimization
        ("4. AI 기반 파라미터 자동 최적화 (PSO)", "• 수많은 모션 파라미터(V_max, A_max, Beta 등)를 사람이 일일이 찾는 것은 비효율적\n• 입자군집최적화(Particle Swarm Optimization) AI 알고리즘 탑재\n• 목표: 최대 속도 한계 내에서 '진동 오차가 0에 도달하는 최단 정착시간' 자동 탐색"),
        # Slide 6: Results
        ("5. 시뮬레이션 결과 및 결론", "• [결과] 제어 미적용(Trapezoidal) 대비 제어 적용(Shaped) 시 잔류진동 정착시간 획기적 단축 확인\n• [효과] 고가의 실제 하드웨어 없이도 최적의 모션 파라미터 사전 튜닝 가능\n• [결론] 장비 셋업 시간 단축 및 소프트웨어 기반의 근본적인 UPH 극대화 솔루션 제시")
    ]
    
    for i, slide in enumerate(prs.slides):
        if i < len(slide_contents):
            title, content = slide_contents[i]
            
            text_shapes = [s for s in slide.shapes if s.has_text_frame]
            title_shape = None
            body_shape = None
            
            for shape in slide.placeholders:
                if shape.placeholder_format.type in [1, 3]: title_shape = shape
                elif shape.placeholder_format.type in [2, 7]: body_shape = shape
                    
            if title_shape and title_shape.has_text_frame:
                title_shape.text_frame.text = title
                for p in title_shape.text_frame.paragraphs: p.font.size = Pt(28)
            elif len(text_shapes) > 0:
                text_shapes[0].text_frame.text = title
                for p in text_shapes[0].text_frame.paragraphs: p.font.size = Pt(28)
                
            if body_shape and body_shape.has_text_frame:
                body_shape.text_frame.text = content
                for p in body_shape.text_frame.paragraphs: p.font.size = Pt(16)
            elif len(text_shapes) > 1:
                text_shapes[1].text_frame.text = content
                for p in text_shapes[1].text_frame.paragraphs: p.font.size = Pt(16)

    prs.save(template_path)

if __name__ == '__main__':
    inject_poster()
    inject_presentation()
    print("Simulation-only content injected.")
