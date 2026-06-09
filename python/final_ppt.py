# -*- coding: utf-8 -*-
"""
final_ppt.py
최종 발표 PPT 생성 스크립트
기존 템플릿의 레이아웃/디자인을 유지하면서 내용만 채워넣습니다.

Template structure (10 slides):
  Slide 1  (제목 슬라이드)     - 표지
  Slide 2  (목차 레이아웃)     - 목차
  Slide 3  (01 분할형)        - 연구 배경 (left text + right table)
  Slide 4  (02 가로형)        - 시스템 하드웨어
  Slide 5  (01 분할형)        - 핵심 제어 알고리즘
  Slide 6  (03 가로형)        - 디지털 트윈
  Slide 7  (01 분할형)        - 검증 결과 (title only → add content)
  Slide 8  (05 가로형)        - 결과분석 및 토의 (title only → add content)
  Slide 9  (06 가로형)        - 느낀점 및 애로사항
  Slide 10 (1_제목 슬라이드)   - 감사합니다 (empty)

We remap these to the user's 10-slide content plan.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from copy import deepcopy
import os

TEMPLATE = r'd:\dongwon\PICKPLACEROBOT\[4학년] (양식)2026년 산학프로젝트 최종보고 발표자료.pptx'
OUTPUT = r'd:\dongwon\PICKPLACEROBOT\[완성본] 2026년 산학프로젝트 최종보고 발표자료.pptx'


def clear_text_frame(tf):
    """Clear all paragraphs in a text frame, keeping the first paragraph."""
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    # Clear the remaining first paragraph
    first_p = tf.paragraphs[0]
    for r in first_p.runs:
        r._r.getparent().remove(r._r)
    first_p.text = ""


def set_text_preserving_format(tf, lines, font_size=None, font_name=None, bold=None, color=None):
    """
    Replace text frame content with given lines.
    Tries to preserve existing run formatting from the first run.
    """
    from lxml import etree
    from pptx.oxml.ns import qn
    
    # Capture formatting from existing first run if available
    ref_font_size = font_size
    ref_font_name = font_name
    ref_font_bold = bold
    ref_font_color = color
    
    if tf.paragraphs and tf.paragraphs[0].runs:
        ref_run = tf.paragraphs[0].runs[0]
        if ref_font_size is None and ref_run.font.size:
            ref_font_size = ref_run.font.size
        if ref_font_name is None and ref_run.font.name:
            ref_font_name = ref_run.font.name
        if ref_font_bold is None and ref_run.font.bold is not None:
            ref_font_bold = ref_run.font.bold
    
    clear_text_frame(tf)
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        run = p.add_run()
        run.text = line
        
        if ref_font_size:
            run.font.size = ref_font_size
        if ref_font_name:
            run.font.name = ref_font_name
        if ref_font_bold is not None:
            run.font.bold = ref_font_bold
        if ref_font_color:
            run.font.color.rgb = ref_font_color


def add_textbox(slide, left, top, width, height, lines, font_size=Pt(14), font_name=None, bold=False, color=None, alignment=None):
    """Add a text box with multiple lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        if font_name:
            run.font.name = font_name
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        if alignment:
            p.alignment = alignment
    
    return txBox


def add_table_to_slide(slide, rows_data, left, top, width, height,
                       header_color=RGBColor(0x1F, 0x49, 0x7D),
                       header_font_color=RGBColor(0xFF, 0xFF, 0xFF)):
    """Add a formatted table to a slide."""
    num_rows = len(rows_data)
    num_cols = len(rows_data[0])
    
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table
    
    # Set column widths proportionally
    col_widths = [Cm(3.5), Cm(4.5), Cm(4.5), Cm(4.5)]
    for ci in range(num_cols):
        if ci < len(col_widths):
            table.columns[ci].width = col_widths[ci]
    
    for ri, row_data in enumerate(rows_data):
        for ci, cell_text in enumerate(row_data):
            cell = table.cell(ri, ci)
            cell.text = ""
            
            # Clear default paragraph and set text
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(cell_text)
            run.font.size = Pt(11)
            
            if ri == 0:
                # Header row
                run.font.bold = True
                run.font.color.rgb = header_font_color
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                run.font.bold = False
                # Highlight the last column (AS-Curve+ZVD) with light blue
                if ci == num_cols - 1 and ri > 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xDB, 0xE5, 0xF1)
            
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return table_shape


def main():
    prs = Presentation(TEMPLATE)
    slides = list(prs.slides)
    
    # =========================================================
    # Slide 1: 표지
    # =========================================================
    slide1 = slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                # Title placeholder
                set_text_preserving_format(shape.text_frame, [
                    "AI 기반 비대칭 S-Curve 궤적 최적화 및",
                    "인풋쉐이핑을 활용한 반도체 픽앤플레이스",
                    "갠트리 진동 억제 시스템"
                ])
            elif shape.name == 'TextBox 3':
                # Subtitle area
                set_text_preserving_format(shape.text_frame, [
                    "2026년 산학프로젝트 최종보고",
                    "디지털 트윈 시뮬레이션 기반 성능 검증"
                ])
            # 지도교수/참여기업 shape — leave as-is (빈칸 유지)
    
    # =========================================================
    # Slide 2: 목차
    # =========================================================
    slide2 = slides[1]
    for shape in slide2.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 1:
            set_text_preserving_format(shape.text_frame, [
                "1. 연구 배경 및 목적",
                "2. 시스템 모델링",
                "3. 모션 프로파일 비교 (AS-Curve & Input Shaping)",
                "4. 인풋쉐이핑 원리 (ZVD)",
                "5. 디지털 트윈 시뮬레이터 구현",
                "6. 실험 결과 — 동일 조건 비교",
                "7. 핵심 결론",
                "8. 한계점 및 향후 과제"
            ])
    
    # =========================================================
    # Slide 3: 연구 배경 및 목적 (01 분할형 — left text + right table)
    # =========================================================
    slide3 = slides[2]
    for shape in slide3.shapes:
        if shape.has_text_frame:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                # Title
                set_text_preserving_format(shape.text_frame, [
                    "1. 연구 배경 및 목적"
                ])
            elif shape.is_placeholder and shape.placeholder_format.idx == 10:
                # Left content
                set_text_preserving_format(shape.text_frame, [
                    "• 반도체 패키징 공정의 Pick & Place 로봇은",
                    "  갠트리(X축)의 고속 이동 시 잔류 진동 발생",
                    "",
                    "• 잔류 진동 → 정착 시간 증가 → UPH 저하",
                    "",
                    "• 목적:",
                    "  AS-Curve 궤적 + ZVD 인풋쉐이핑으로",
                    "  잔류 진동을 근본 억제하여 UPH 극대화"
                ])
    
    # =========================================================
    # Slide 4: 시스템 모델링 (02 가로형)
    # =========================================================
    slide4 = slides[3]
    for shape in slide4.shapes:
        if shape.has_text_frame:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                set_text_preserving_format(shape.text_frame, [
                    "2. 시스템 모델링: 2-Mass Perfect Servo"
                ])
            elif shape.name == '내용 개체 틀 1':
                set_text_preserving_format(shape.text_frame, [
                    "• 2-Mass Perfect Servo 모델",
                    "  - 모터(캐리지): 명령 궤적을 완벽히 추종",
                    "  - 로드(엔드이펙터): 스프링-댐퍼로 캐리지에 연결",
                    "",
                    "• 물리 방정식",
                    "  F = k·(x₁ - x₂) + c·(v₁ - v₂)",
                    "  a₂ = F / m₂",
                    "",
                    "• 수치 적분: RK4 (10× substep, dt = 0.1 ms)",
                    "",
                    "• 설계 파라미터",
                    "  - 고유진동수 fn = 10 Hz",
                    "  - 감쇠비 ζ = 0.05"
                ])
    
    # =========================================================
    # Slide 5: 모션 프로파일 비교 (01 분할형)
    # =========================================================
    slide5 = slides[4]
    for shape in slide5.shapes:
        if shape.has_text_frame:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                set_text_preserving_format(shape.text_frame, [
                    "3. 모션 프로파일 비교"
                ])
            elif shape.is_placeholder and shape.placeholder_format.idx == 11:
                set_text_preserving_format(shape.text_frame, [
                    "① Trapezoidal",
                    "  • 가속도가 순간 변화 → 저크 무제한",
                    "  • 최대 진동 유발 (비교 기준)",
                    "",
                    "② AS-Curve (비대칭 S-Curve)",
                    "  • 7-segment 비대칭 저크 제한",
                    "  • 기구부 충격 완화",
                    "",
                    "③ AS-Curve + ZVD",
                    "  • 인풋쉐이핑으로 잔류 진동 근본 상쇄",
                    "  • 정착 시간 획기적 단축"
                ])
    
    # =========================================================
    # Slide 6: 디지털 트윈 시뮬레이터 (03 가로형)
    # =========================================================
    slide6 = slides[5]
    for shape in slide6.shapes:
        if shape.has_text_frame:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                set_text_preserving_format(shape.text_frame, [
                    "4. 인풋쉐이핑 원리: ZVD 쉐이퍼"
                ])
            elif shape.name == '내용 개체 틀 1':
                set_text_preserving_format(shape.text_frame, [
                    "• ZVD(Zero Vibration and Derivative) 쉐이퍼",
                    "  - 3개 임펄스를 고유진동 반주기 간격으로 배치",
                    "",
                    "• 동작 원리",
                    "  - 파괴적 간섭(Destructive Interference)으로",
                    "    잔류 진동을 이론적으로 0으로 억제",
                    "",
                    "• 쉐이퍼 지연: ~100 ms",
                    "  - 고속 이동 시간(~95 ms) 대비 매우 짧음",
                    "  - 이 지연 투자로 2.5초의 진동 대기 제거"
                ])
    
    # =========================================================
    # Slide 7: 검증 결과 (01 분할형 — title only, need to add content)
    # This slide originally has only the title. We'll add content.
    # =========================================================
    slide7 = slides[6]
    for shape in slide7.shapes:
        if shape.has_text_frame:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                set_text_preserving_format(shape.text_frame, [
                    "5. 디지털 트윈 시뮬레이터 구현"
                ])
    
    # Add content text box to slide 7
    add_textbox(slide7,
                left=Cm(2.2), top=Cm(5.6), width=Cm(24), height=Cm(8),
                lines=[
                    "• Three.js 기반 3D 실시간 시각화",
                    "  - Shaped(제어 적용) vs Unshaped(고스트) 로봇 동시 비교",
                    "",
                    "• Chart.js 기반 실시간 분석 차트",
                    "  - 위치 / 속도 / 가속도 / 잔류진동 4채널",
                    "",
                    "• AI PSO(Particle Swarm Optimization) 자동 튜닝 탑재",
                    "  - V_max, A_max, Beta, Gamma 최적 파라미터 자동 탐색",
                    "",
                    "• 웹 브라우저에서 즉시 구동 가능한 경량 시뮬레이터"
                ],
                font_size=Pt(14))
    
    # =========================================================
    # Slide 8: 결과분석 (05 가로형 — title only, need to add content)
    # This becomes the key results slide with the comparison table
    # =========================================================
    slide8 = slides[7]
    for shape in slide8.shapes:
        if shape.has_text_frame:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                set_text_preserving_format(shape.text_frame, [
                    "6. 실험 결과 — 동일 조건 비교 ★"
                ])
    
    # Add the comparison table
    table_data = [
        ["지표", "Trapezoidal", "AS-Curve", "AS-Curve + ZVD"],
        ["이동 시간", "0.1000 s", "0.0950 s", "0.0950 s"],
        ["잔류 진동", "17.47 mm", "24.66 mm", "4.29 mm"],
        ["정착 시간", "2.660 s", "2.815 s", "0.197 s"],
        ["UPH", "676", "639", "9,137"],
    ]
    
    add_table_to_slide(slide8, table_data,
                       left=Cm(2.5), top=Cm(4.8),
                       width=Cm(23), height=Cm(6))
    
    # Add conditions note below
    add_textbox(slide8,
                left=Cm(2.5), top=Cm(11.5), width=Cm(23), height=Cm(2),
                lines=[
                    "조건: V=1,000 mm/s, A=20,000 mm/s², fn=10 Hz, ζ=0.05, 이동거리=50 mm"
                ],
                font_size=Pt(11),
                color=RGBColor(0x66, 0x66, 0x66))
    
    # =========================================================
    # Slide 9: 핵심 결론 + 한계점 (06 가로형)
    # =========================================================
    slide9 = slides[8]
    for shape in slide9.shapes:
        if shape.has_text_frame:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                set_text_preserving_format(shape.text_frame, [
                    "7. 핵심 결론 & 8. 한계점 및 향후 과제"
                ])
            elif shape.name == '내용 개체 틀 1':
                set_text_preserving_format(shape.text_frame, [
                    "【핵심 결론】",
                    "  • 동일 조건 대비 잔류 진동 75.4% 감소",
                    "  • 정착 시간 92.6% 단축 (2.660s → 0.197s)",
                    "  • UPH 1,252% 향상 (676 → 9,137)",
                    "  • ZVD의 100ms 지연 투자 → 2.5초 진동 대기 제거",
                    "",
                    "【한계점 및 향후 과제】",
                    "  • Perfect Servo 가정 → 실제 서보 추종 오차 미반영",
                    "  • 단일 축(X) 모델 → 다축 커플링 미고려",
                    "  • 향후: 실물 HW 검증 (STM32 + NEMA 스텝모터 + IMU 센서)"
                ])
    
    # =========================================================
    # Slide 10: 감사합니다 / Q&A (1_제목 슬라이드 — empty)
    # =========================================================
    slide10 = slides[9]
    
    # Add "감사합니다" centered
    add_textbox(slide10,
                left=Cm(3), top=Cm(5), width=Cm(22), height=Cm(4),
                lines=["감사합니다"],
                font_size=Pt(44),
                bold=True,
                color=RGBColor(0x1F, 0x49, 0x7D),
                alignment=PP_ALIGN.CENTER)
    
    # Add Q&A subtitle
    add_textbox(slide10,
                left=Cm(3), top=Cm(9.5), width=Cm(22), height=Cm(2),
                lines=["Q & A"],
                font_size=Pt(28),
                bold=False,
                color=RGBColor(0x66, 0x66, 0x66),
                alignment=PP_ALIGN.CENTER)
    
    # =========================================================
    # Save
    # =========================================================
    prs.save(OUTPUT)
    print(f"✅ 최종 발표자료 생성 완료: {OUTPUT}")
    print(f"   슬라이드 수: {len(prs.slides)}")


if __name__ == '__main__':
    main()
