"""
Rebuild final presentation PPT with research charts and diagrams embedded.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

TEMPLATE = r'd:\dongwon\PICKPLACEROBOT\[4학년] (양식)2026년 산학프로젝트 최종보고 발표자료.pptx'
OUTPUT = r'd:\dongwon\PICKPLACEROBOT\[완성본] 2026년 산학프로젝트 최종보고 발표자료.pptx'
CHARTS = r'd:\dongwon\PICKPLACEROBOT\charts'
DIAGRAM = r'C:\Users\41110\.gemini\antigravity\brain\39415a21-1eec-4c73-a2d3-203207d31a90\system_diagram_1781031157804.png'

prs = Presentation(TEMPLATE)
slide_width = prs.slide_width
slide_height = prs.slide_height

def clear_slide_texts(slide):
    """Clear all text from a slide's text frames"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""

def add_textbox(slide, left, top, width, height, text, font_size=11, bold=False, color=RGBColor(0x1f, 0x29, 0x37), align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=11, color=RGBColor(0x1f, 0x29, 0x37)):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, sz) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz if sz else font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return txBox

# Slide dimensions in EMU
SW = int(slide_width)
SH = int(slide_height)

# Helper: common margins
ML = int(SW * 0.06)  # left margin
MT = int(SH * 0.18)  # top margin (below title area)
CW = int(SW * 0.88)  # content width

slides = list(prs.slides)

# ================================================
# Slide 1: Title (keep template, just fill in)
# ================================================
s = slides[0]
for shape in s.shapes:
    if shape.has_text_frame:
        full = shape.text_frame.text.strip()
        if '프로젝트명' in full or '제목' in full or len(full) < 5:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.text = ""

add_textbox(s, ML, int(SH*0.30), CW, int(SH*0.20),
    "AI 기반 비대칭 S-Curve 궤적 최적화 및 인풋쉐이핑을 활용한\n반도체 픽앤플레이스 갠트리 진동 억제 시스템",
    font_size=20, bold=True, color=RGBColor(0x1a, 0x1a, 0x2e))

add_textbox(s, ML, int(SH*0.53), CW, int(SH*0.08),
    "디지털 트윈 시뮬레이션 기반 성능 검증",
    font_size=14, bold=False, color=RGBColor(0x4b, 0x55, 0x63))

# ================================================
# Slide 2: TOC (keep as is if already filled)
# ================================================

# ================================================
# Slide 3: Background + Purpose
# ================================================
s = slides[2]
add_multi_text(s, ML, MT, int(CW*0.5), int(SH*0.7), [
    ("1. 연구 배경 및 목적", True, 16),
    ("", False, 6),
    ("■ 연구 배경", True, 12),
    ("• 반도체 패키징의 Pick & Place 로봇은 X-Gantry의", False, 10),
    ("  고속 이동 시 잔류 진동이 발생", False, 10),
    ("• 진동이 정착될 때까지 대기 → UPH 저하의 직접 원인", False, 10),
    ("• 기존 사다리꼴 프로파일: 무한 저크 → 최대 진동 유발", False, 10),
    ("", False, 6),
    ("■ 연구 목적", True, 12),
    ("• AS-Curve 궤적 + ZVD 인풋쉐이핑으로", False, 10),
    ("  잔류 진동을 근본 억제하여 UPH 극대화", False, 10),
    ("• 디지털 트윈 시뮬레이션으로 정량적 검증", False, 10),
])

# ================================================
# Slide 4: System Modeling + Diagram Image
# ================================================
s = slides[3]
add_multi_text(s, ML, MT, int(CW*0.45), int(SH*0.7), [
    ("2. 시스템 모델링", True, 16),
    ("", False, 6),
    ("■ 2-Mass Perfect Servo Model", True, 12),
    ("• Motor(m₁): 위치 명령 완벽 추종", False, 10),
    ("• Load(m₂): 스프링-댐퍼로 연결", False, 10),
    ("", False, 4),
    ("F = k·(x₁-x₂) + c·(ẋ₁-ẋ₂)", True, 11),
    ("a₂ = F / m₂", True, 11),
    ("", False, 4),
    ("■ 수치 해석", True, 12),
    ("• RK4 적분 (10x substep, dt=0.1ms)", False, 10),
    ("• fn=10Hz, ζ=0.05", False, 10),
    ("• Payload 시 m₂ 2배 → fn 자연 하락", False, 10),
])

if os.path.exists(DIAGRAM):
    s.shapes.add_picture(DIAGRAM, Emu(int(SW*0.50)), Emu(int(SH*0.22)), Emu(int(SW*0.45)), Emu(int(SH*0.55)))

# ================================================
# Slide 5: Motion Profile Comparison + Chart 5
# ================================================
s = slides[4]
add_multi_text(s, ML, MT, int(CW*0.40), int(SH*0.65), [
    ("3. 모션 프로파일 비교", True, 16),
    ("", False, 6),
    ("■ Trapezoidal", True, 11),
    ("  순간 가속도 변화 → 최대 진동", False, 10),
    ("■ AS-Curve (비대칭 S-Curve)", True, 11),
    ("  7-segment 저크 제한 → 충격 완화", False, 10),
    ("■ AS-Curve + ZVD", True, 11),
    ("  인풋쉐이핑으로 진동 근본 상쇄", False, 10),
])

chart5 = os.path.join(CHARTS, 'chart5_3profile_bar.png')
if os.path.exists(chart5):
    s.shapes.add_picture(chart5, Emu(int(SW*0.42)), Emu(int(SH*0.25)), Emu(int(SW*0.55)), Emu(int(SH*0.50)))

# ================================================
# Slide 6: Input Shaping Principle + Sensitivity
# ================================================
s = slides[5]
add_multi_text(s, ML, MT, int(CW*0.42), int(SH*0.65), [
    ("4. 인풋쉐이핑 원리", True, 16),
    ("", False, 6),
    ("■ ZVD (Zero Vibration & Derivative)", True, 12),
    ("• 3개 임펄스를 반주기 간격으로 배치", False, 10),
    ("• 파괴적 간섭 → 잔류 진동 이론적 0", False, 10),
    ("• 지연: ~100ms (이동 시간 대비 매우 짧음)", False, 10),
    ("", False, 4),
    ("■ ZVD vs ZV vs EI 강건성 비교", True, 12),
    ("• ZVD: fn 오차 ±15%까지 진동 5% 이하", False, 10),
    ("• EI: 가장 넓은 강건성 대역폭", False, 10),
])

chart3 = os.path.join(CHARTS, 'chart3_sensitivity.png')
if os.path.exists(chart3):
    s.shapes.add_picture(chart3, Emu(int(SW*0.44)), Emu(int(SH*0.20)), Emu(int(SW*0.52)), Emu(int(SH*0.52)))

# ================================================
# Slide 7: Simulator + Vibration Waveform
# ================================================
s = slides[6]
add_multi_text(s, ML, MT, int(CW*0.40), int(SH*0.65), [
    ("5. 디지털 트윈 시뮬레이터", True, 16),
    ("", False, 6),
    ("■ 구현 기술", True, 12),
    ("• Three.js 3D 실시간 렌더링", False, 10),
    ("• Shaped vs Unshaped 동시 비교", False, 10),
    ("• Chart.js 4채널 실시간 차트", False, 10),
    ("• AI PSO 자동 파라미터 튜닝", False, 10),
    ("", False, 4),
    ("■ 우측: 잔류 진동 파형 비교", True, 12),
    ("  (Trapezoidal vs AS-Curve+ZVD)", False, 10),
])

chart1 = os.path.join(CHARTS, 'chart1_vibration_waveform.png')
if os.path.exists(chart1):
    s.shapes.add_picture(chart1, Emu(int(SW*0.42)), Emu(int(SH*0.22)), Emu(int(SW*0.55)), Emu(int(SH*0.50)))

# ================================================
# Slide 8: Results Table + Speed Sweep
# ================================================
s = slides[7]
# Add table
from pptx.util import Inches
rows, cols = 5, 4
tbl = s.shapes.add_table(rows, cols, Emu(ML), Emu(int(SH*0.20)), Emu(int(CW*0.50)), Emu(int(SH*0.38))).table

headers = ['지표', 'Trapezoidal', 'AS-Curve', 'AS-Curve+ZVD']
data = [
    ['이동 시간', '0.1000 s', '0.0950 s', '0.0950 s'],
    ['잔류 진동', '17.47 mm', '24.66 mm', '4.29 mm'],
    ['정착 시간', '2.660 s', '2.815 s', '0.197 s'],
    ['UPH', '676', '639', '9,137'],
]

for c, h in enumerate(headers):
    cell = tbl.cell(0, c)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x1e, 0x3a, 0x5f)

for r, row_data in enumerate(data):
    for c, val in enumerate(row_data):
        cell = tbl.cell(r+1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.alignment = PP_ALIGN.CENTER
            if c == 3:  # Highlight ZVD column
                p.font.bold = True
                p.font.color.rgb = RGBColor(0x1a, 0x56, 0xdb)

add_textbox(s, ML, int(SH*0.60), int(CW*0.50), int(SH*0.08),
    "조건: V=1,000 mm/s, A=20,000 mm/s², fn=10Hz, ζ=0.05, 거리=50mm",
    font_size=8, color=RGBColor(0x6b, 0x72, 0x80))

# Speed sweep chart on right half
chart2 = os.path.join(CHARTS, 'chart2_uph_vs_speed.png')
if os.path.exists(chart2):
    s.shapes.add_picture(chart2, Emu(int(SW*0.52)), Emu(int(SH*0.20)), Emu(int(SW*0.44)), Emu(int(SH*0.42)))

# Distance chart below
chart4 = os.path.join(CHARTS, 'chart4_settle_vs_distance.png')
if os.path.exists(chart4):
    s.shapes.add_picture(chart4, Emu(int(SW*0.52)), Emu(int(SH*0.62)), Emu(int(SW*0.44)), Emu(int(SH*0.33)))

add_textbox(s, ML, int(SH*0.15), int(CW), int(SH*0.05),
    "6. 실험 결과 — 동일 조건 비교 + 파라미터 스윕",
    font_size=15, bold=True, color=RGBColor(0x1e, 0x3a, 0x5f))

# ================================================
# Slide 9: Conclusion + Economic Impact
# ================================================
s = slides[8]
add_multi_text(s, ML, MT, int(CW*0.42), int(SH*0.7), [
    ("7. 핵심 결론", True, 16),
    ("", False, 6),
    ("■ 정량적 성과 (Trapezoidal 대비)", True, 12),
    ("• 잔류 진동 75.4% 감소", False, 11),
    ("• 정착 시간 92.6% 단축 (2.66s → 0.20s)", False, 11),
    ("• UPH 1,252% 향상 (676 → 9,137)", False, 11),
    ("", False, 6),
    ("■ 한계점 및 향후 과제", True, 12),
    ("• Perfect Servo 가정 (추종 오차 미반영)", False, 10),
    ("• 단일 축(X) 모델 → 다축 커플링 미고려", False, 10),
    ("• 향후: STM32 + IMU 실물 검증", False, 10),
])

chart6 = os.path.join(CHARTS, 'chart6_economic.png')
if os.path.exists(chart6):
    s.shapes.add_picture(chart6, Emu(int(SW*0.44)), Emu(int(SH*0.22)), Emu(int(SW*0.52)), Emu(int(SH*0.48)))

# ================================================
# Slide 10: Thank you
# ================================================
s = slides[9]
add_textbox(s, int(SW*0.25), int(SH*0.35), int(SW*0.50), int(SH*0.15),
    "감사합니다",
    font_size=32, bold=True, color=RGBColor(0x1e, 0x3a, 0x5f), align=PP_ALIGN.CENTER)
add_textbox(s, int(SW*0.25), int(SH*0.52), int(SW*0.50), int(SH*0.08),
    "Q & A",
    font_size=18, bold=False, color=RGBColor(0x6b, 0x72, 0x80), align=PP_ALIGN.CENTER)

prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
print("Images embedded:")
for f in ['chart1_vibration_waveform.png', 'chart2_uph_vs_speed.png', 'chart3_sensitivity.png',
          'chart4_settle_vs_distance.png', 'chart5_3profile_bar.png', 'chart6_economic.png']:
    fp = os.path.join(CHARTS, f)
    print(f"  {'OK' if os.path.exists(fp) else 'MISSING'} {f}")
print(f"  {'OK' if os.path.exists(DIAGRAM) else 'MISSING'} system_diagram.png")
