import os
from pptx import Presentation
from pptx.util import Pt

def inject_poster():
    template_path = '[4학년] (양식)2026년 산학프로젝트 포스터_260512.ver.pptx'
    if not os.path.exists(template_path):
        return
        
    prs = Presentation(template_path)
    slide = prs.slides[0]
    
    # 포스터의 모든 텍스트 프레임을 찾아서 제목과 내용을 적절히 덮어씁니다.
    # 하지만 보통 포스터 양식은 한 장의 큰 이미지 위에 빈 텍스트 상자가 있거나, 
    # 여러 개의 텍스트 상자가 있습니다.
    # 사용자의 요청에 따라 강제로 텍스트를 가장 큰 텍스트 상자나 제목에 밀어넣습니다.
    
    title_text = "반도체 Pick & Place 장비의 UPH 극대화를 위한 잔류 진동 억제 궤적 설계 및 디지털 트윈"
    content_text = """[연구 배경 및 목적]
- 반도체 패키징 공정 생산성(UPH) 향상을 위해 고속 이송 시 로봇 끝단에 잔류 진동 발생.
- 정착 시간을 최소화하면서 가장 빠른 이송이 가능한 궤적 최적화 목표.

[핵심 기술]
1. 비대칭 S-Curve (AS-Curve): 가속은 빠르게, 감속은 부드럽게 설계하여 물리적 충격 최소화
2. Input Shaping: 시스템의 고유진동수(fn)를 분석해 역위상 임펄스를 인가하여 진동 상쇄
3. AI 최적화 (PSO): 최소 정착 시간을 갖는 파라미터(V_max, A_max, Beta 등) 자동 도출

[디지털 트윈 및 결과]
- 실제 동역학 모델을 웹 환경의 3D 시뮬레이션으로 구현. 
- 제어 미적용(Baseline) 대비 제어 적용(Shaped) 시 잔류 진동 오차가 0으로 수렴함을 입증함."""
    
    # 텍스트가 있는 프레임들 중, 제목과 본문으로 추정되는 것들을 변경
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    if text_shapes:
        # 첫 번째 도형을 제목으로 간주
        text_shapes[0].text_frame.text = title_text
        if len(text_shapes) > 1:
            # 두 번째 도형을 본문으로 간주
            text_shapes[1].text_frame.text = content_text
        else:
            # 도형이 하나밖에 없으면 강제 추가
            txBox = slide.shapes.add_textbox(Pt(50), Pt(150), prs.slide_width - Pt(100), prs.slide_height - Pt(200))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.text = content_text
    
    prs.save(template_path)


def inject_presentation():
    template_path = '[4학년] (양식)2026년 산학프로젝트 최종보고 발표자료.pptx'
    if not os.path.exists(template_path):
        return
        
    prs = Presentation(template_path)
    
    slide_contents = [
        # Slide 0: Title
        ("반도체 Pick & Place 장비의 UPH 극대화를 위한 잔류 진동 억제 궤적 설계", "2026년 산학프로젝트 최종보고"),
        # Slide 1: 목차 (Skip or fill with generic contents)
        ("목차", "1. 연구 배경\n2. AS-Curve 궤적 설계\n3. Input Shaping 제어\n4. 디지털 트윈 및 AI 최적화\n5. 결론 및 기대효과"),
        # Slide 2: Background
        ("1. 연구 배경: 반도체 장비의 딜레마", "• 생산성(UPH) 향상을 위해 모터 속도 증가 필수\n• 고속 이동 시 로봇 팔 끝단에 심각한 '잔류 진동(Residual Vibration)' 발생\n• 진동이 멈출 때까지 기다려야 하는 '정착 시간(Settling Time)'이 길어지면 전체 공정 시간 증가"),
        # Slide 3: Solution 1
        ("2. 핵심 솔루션 1: 비대칭 S-Curve 궤적", "• 기존 사다리꼴(Trapezoidal) 궤적: 가속도 불연속으로 큰 저크(Jerk) 발생\n• AS-Curve (비대칭 S-Curve) 적용:\n  - 가속 구간: 빠르게 도달\n  - 감속 구간: 부드럽게 정지 (Beta, Gamma 파라미터 최적화)\n• 물리적 충격량 최소화 입증"),
        # Slide 4: Solution 2
        ("3. 핵심 솔루션 2: Input Shaping 제어", "• 파동의 상쇄 간섭 원리 활용\n• 첫 번째 충격으로 발생한 진동을, 두 번째 충격(임펄스)으로 정확히 상쇄(Cancel out)\n• 로봇 팔의 고유진동수(fn)와 감쇠비(zeta)를 계산하여 맞춤형 쉐이퍼 설계"),
        # Slide 5: Digital Twin
        ("4. 디지털 트윈 및 AI 자동 튜닝 개발", "• 이론적 계산을 넘어 실제 3D 시뮬레이션 환경(Digital Twin) 구현\n• 4가지 고해상도 실시간 분석 차트 탑재\n• AI 최적화 (Particle Swarm Optimization):\n  - 사람이 찾기 힘든 최적의 변수를 백그라운드에서 자동 탐색"),
        # Slide 6: Results
        ("5. 결과 분석 및 기대효과", "• [결과] 제어 미적용 시 대비 정착 시간 대폭 감소 및 진동 제거 확인\n• [효과] 오프라인 환경에서 모터 파라미터 사전 튜닝 가능\n• [결론] 장비 셋업 시간 단축 및 실질적 UPH 향상에 기여")
    ]
    
    # 템플릿의 슬라이드를 순회하며 내용을 채워넣음
    for i, slide in enumerate(prs.slides):
        if i < len(slide_contents):
            title, content = slide_contents[i]
            
            # 슬라이드의 텍스트 프레임들을 찾음
            text_shapes = [s for s in slide.shapes if s.has_text_frame]
            
            # 플레이스홀더를 사용하여 제목/본문 구분 시도
            title_shape = None
            body_shape = None
            
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1 or shape.placeholder_format.type == 3: # TITLE or CENTER_TITLE
                    title_shape = shape
                elif shape.placeholder_format.type == 2 or shape.placeholder_format.type == 7: # BODY or OBJECT
                    body_shape = shape
                    
            if title_shape and title_shape.has_text_frame:
                title_shape.text_frame.text = title
            elif len(text_shapes) > 0:
                text_shapes[0].text_frame.text = title
                
            if body_shape and body_shape.has_text_frame:
                body_shape.text_frame.text = content
            elif len(text_shapes) > 1:
                text_shapes[1].text_frame.text = content
                
            # 만약 슬라이드에 아무 텍스트 프레임이 없다면 강제로 추가
            if not title_shape and not body_shape and len(text_shapes) == 0:
                txBox = slide.shapes.add_textbox(Pt(50), Pt(50), prs.slide_width - Pt(100), Pt(50))
                txBox.text_frame.text = title
                txBox2 = slide.shapes.add_textbox(Pt(50), Pt(120), prs.slide_width - Pt(100), prs.slide_height - Pt(150))
                txBox2.text_frame.text = content

    prs.save(template_path)

if __name__ == '__main__':
    inject_poster()
    inject_presentation()
    print("Direct injection complete.")
