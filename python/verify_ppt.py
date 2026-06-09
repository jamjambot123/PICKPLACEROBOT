# -*- coding: utf-8 -*-
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

prs = Presentation(r'd:\dongwon\PICKPLACEROBOT\[완성본] 2026년 산학프로젝트 최종보고 발표자료.pptx')
print(f'Total slides: {len(prs.slides)}')
for i, slide in enumerate(prs.slides):
    print(f'\n=== Slide {i+1} (layout: {slide.slide_layout.name}) ===')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pi, p in enumerate(shape.text_frame.paragraphs):
                txt = p.text.strip()
                if txt:
                    print(f'  [{shape.name}] {txt[:120]}')
        if shape.has_table:
            table = shape.table
            ncols = len(table.columns)
            print(f'  TABLE ({len(table.rows)}x{ncols}):')
            for ri, row in enumerate(table.rows):
                cells = []
                for ci in range(ncols):
                    cells.append(row.cells[ci].text)
                sep = ' | '
                print(f'    Row {ri}: {sep.join(cells)}')
