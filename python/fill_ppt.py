import os
from pptx import Presentation
from pptx.util import Pt

def fill_poster():
    template_path = '[4학년] (양식)2026년 산학프로젝트 포스터_260512.ver.pptx'
    if not os.path.exists(template_path):
        return
        
    prs = Presentation(template_path)
    
    # 포스터는 보통 1장이므로, 첫 번째 슬라이드에 텍스트 상자를 추가하여 내용을 넣어줍니다.
    slide = prs.slides[0]
    
    # 내용 요약 텍스트 박스 추가 (배경을 가리지 않게 상단이나 여백에 추가)
    left = top = width = height = 0  # 임시
    txBox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    tf.text = "[AI 자동 완성 초안 - 디자인에 맞게 옮겨 적어주세요!]\n\n"
    
    p = tf.add_paragraph()
    p.text = "[연구 배경 및 목적]\n반도체 패키징 공정 생산성(UPH) 향상을 위해 고속 이송 시 로봇 끝단에 잔류 진동 발생. 정착 시간을 최소화하면서 가장 빠른 이송이 가능한 궤적 최적화 목표.\n\n"
    
    p = tf.add_paragraph()
    p.text = "[핵심 기술]\n1. 비대칭 S-Curve (AS-Curve): 가속은 빠르게, 감속은 부드럽게 설계하여 물리적 충격 최소화\n2. Input Shaping: 시스템의 고유진동수(fn)를 분석해 역위상 임펄스를 인가하여 진동 상쇄\n3. AI 최적화 (PSO): 최소 정착 시간을 갖는 파라미터(V_max, A_max, Beta 등) 자동 도출\n\n"
    
    p = tf.add_paragraph()
    p.text = "[디지털 트윈 및 결과]\n실제 동역학 모델을 웹 환경의 3D 시뮬레이션으로 구현. 제어 미적용(Baseline) 대비 제어 적용(Shaped) 시 잔류 진동 오차가 0으로 수렴함을 입증함."
    
    # 글자 크기 조정
    for para in tf.paragraphs:
        if para.runs:
            para.font.size = Pt(24)
            para.font.bold = True
            
    prs.save('[완성본] 2026년 산학프로젝트 포스터.pptx')


def add_slide(prs, title, content):
    # 보통 레이아웃 1번이 제목+내용
    try:
        slide_layout = prs.slide_layouts[1]
    except:
        slide_layout = prs.slide_layouts[0]
        
    slide = prs.slides.add_slide(slide_layout)
    
    try:
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
    except:
        # 플레이스홀더가 없는 양식일 경우 텍스트 상자 강제 생성
        txBox = slide.shapes.add_textbox(Pt(50), Pt(100), prs.slide_width - Pt(100), prs.slide_height - Pt(150))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = f"{title}\n\n{content}"
    return slide

def fill_presentation():
    template_path = '[4학년] (양식)2026년 산학프로젝트 최종보고 발표자료.pptx'
    if not os.path.exists(template_path):
        return
        
    prs = Presentation(template_path)
    
    add_slide(prs, "1. 연구 배경: 반도체 장비의 딜레마", 
              "• 생산성(UPH) 향상을 위해 모터 속도 증가 필수\n• 고속 이동 시 로봇 팔 끝단에 심각한 '잔류 진동(Residual Vibration)' 발생\n• 진동이 멈출 때까지 기다려야 하는 '정착 시간(Settling Time)'이 길어지면 전체 공정 시간 증가")
              
    add_slide(prs, "2. 핵심 솔루션 1: 비대칭 S-Curve 궤적", 
              "• 기존 사다리꼴(Trapezoidal) 궤적: 가속도 불연속으로 큰 저크(Jerk) 발생\n• AS-Curve (비대칭 S-Curve) 적용:\n  - 가속 구간: 빠르게 도달\n  - 감속 구간: 부드럽게 정지 (Beta, Gamma 파라미터 최적화)\n• 물리적 충격량 최소화 입증")
              
    add_slide(prs, "3. 핵심 솔루션 2: Input Shaping 제어", 
              "• 파동의 상쇄 간섭 원리 활용\n• 첫 번째 충격으로 발생한 진동을, 두 번째 충격(임펄스)으로 정확히 상쇄(Cancel out)\n• 로봇 팔의 고유진동수(fn)와 감쇠비(zeta)를 계산하여 맞춤형 쉐이퍼 설계")

    add_slide(prs, "4. 디지털 트윈 및 AI 자동 튜닝 개발", 
              "• 이론적 계산을 넘어 실제 3D 시뮬레이션 환경(Digital Twin) 구현\n• 4가지 고해상도 실시간 분석 차트 탑재\n• AI 최적화 (Particle Swarm Optimization):\n  - 사람이 찾기 힘든 최적의 변수를 백그라운드에서 자동 탐색")
              
    add_slide(prs, "5. 결과 분석 및 기대효과", 
              "• [결과] 제어 미적용 시 대비 정착 시간 대폭 감소 및 진동 제거 확인\n• [효과] 오프라인 환경에서 모터 파라미터 사전 튜닝 가능\n• [결론] 장비 셋업 시간 단축 및 실질적 UPH 향상에 기여")

    prs.save('[완성본] 2026년 산학프로젝트 최종보고 발표자료.pptx')

if __name__ == '__main__':
    fill_poster()
    fill_presentation()
    print("Templates successfully modified and saved as [완성본].")
