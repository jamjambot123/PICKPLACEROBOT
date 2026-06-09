import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_poster():
    prs = Presentation()
    # A0 size approx for poster (portrait or landscape, usually portrait 33x46 inches)
    prs.slide_width = Inches(33)
    prs.slide_height = Inches(46)
    
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(31), Inches(3))
    tf = txBox.text_frame
    tf.text = "반도체 Pick & Place 장비의 UPH 극대화를 위한 잔류 진동 억제 궤적 설계 및 디지털 트윈"
    tf.paragraphs[0].font.size = Pt(80)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Section 1: Background
    txBox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(14.5), Inches(10))
    tf = txBox.text_frame
    tf.text = "1. 연구 배경 및 목적 (Background & Objective)"
    tf.paragraphs[0].font.size = Pt(60)
    tf.paragraphs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = "- [문제 제기] 반도체 패키징 공정에서 생산성(UPH) 향상을 위해 고속 이송 시 끝단(End-Effector)에 잔류 진동 발생\n- [목표] 정착 시간(Settling Time)을 최소화하면서도 가장 빠른 이송이 가능한 궤적 최적화"
    p.font.size = Pt(40)
    
    # Section 2: Core Tech
    txBox = slide.shapes.add_textbox(Inches(1), Inches(16), Inches(14.5), Inches(12))
    tf = txBox.text_frame
    tf.text = "2. 핵심 알고리즘 (Core Technologies)"
    tf.paragraphs[0].font.size = Pt(60)
    tf.paragraphs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = "- 비대칭 S-Curve (AS-Curve): 가속은 빠르게, 감속은 부드럽게 설계하여 물리적 충격 최소화\n- Input Shaping (입력 성형): 시스템의 고유진동수를 분석해 역위상 임펄스를 인가, 진동 상쇄 (ZV, ZVD 적용)\n- AI 최적화 (PSO): 백그라운드 최적화로 최소 정착 시간을 갖는 파라미터(V_max, A_max, Beta 등) 자동 도출"
    p.font.size = Pt(40)
    
    # Section 3: Digital Twin & Simulation
    txBox = slide.shapes.add_textbox(Inches(16.5), Inches(5), Inches(15.5), Inches(15))
    tf = txBox.text_frame
    tf.text = "3. 디지털 트윈 시뮬레이션 개발"
    tf.paragraphs[0].font.size = Pt(60)
    tf.paragraphs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = "- [웹 기반 3D 환경 구축] 실제 물리 동역학 방정식(ODE)을 웹에 이식하여 실시간 교차 검증\n- [여기에 웹 시뮬레이터 3D 풀스크린 캡처 이미지 삽입 요망!]"
    p.font.size = Pt(40)
    
    # Section 4: Results
    txBox = slide.shapes.add_textbox(Inches(16.5), Inches(22), Inches(15.5), Inches(15))
    tf = txBox.text_frame
    tf.text = "4. 성능 평가 및 결과 (Results)"
    tf.paragraphs[0].font.size = Pt(60)
    tf.paragraphs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = "- 제어 미적용(Baseline) 대비 제어 적용(Shaped) 시 잔류 진동 수렴 시간 획기적 단축\n- [여기에 웹 시뮬레이터 우측의 '잔류 진동 오차' 차트 캡처 이미지 삽입 요망!]"
    p.font.size = Pt(40)
    
    prs.save('[4학년] (작성가이드)2026년 산학프로젝트 포스터_260512.ver.pptx')

def add_slide(prs, title, content, notes=""):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = title
    tf = body_shape.text_frame
    tf.text = content
    
    # Add notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes
    return slide

def create_presentation():
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "반도체 Pick & Place 장비의 UPH 극대화를 위한\n잔류 진동 억제 궤적 설계 및 디지털 트윈"
    slide.shapes.placeholders[1].text = "2026년 산학프로젝트 최종보고\n발표자: OOO"
    
    # Slide 2: 배경
    add_slide(prs, "1. 연구 배경: 반도체 장비의 딜레마", 
              "• 생산성(UPH) 향상을 위해 모터 속도 증가 필수\n• 하지만 고속 이동 시 로봇 팔 끝단에 심각한 '잔류 진동(Residual Vibration)' 발생\n• 진동이 멈출 때까지 기다려야 하는 '정착 시간(Settling Time)'이 길어지면 오히려 전체 공정 시간 증가",
              notes="[발표 멘트] 안녕하십니까, 저희는 반도체 장비의 생산성과 진동이라는 두 마리 토끼를 잡기 위한 연구를 진행했습니다. 빨리 움직이면 흔들리고, 흔들리면 기다려야 하는 딜레마를 해결하고자 했습니다.")
              
    # Slide 3: 핵심 솔루션 1
    add_slide(prs, "2. 솔루션 1: 비대칭 S-Curve 궤적", 
              "• 기존 사다리꼴(Trapezoidal) 궤적: 가속도 불연속으로 큰 저크(Jerk) 발생\n• AS-Curve (비대칭 S-Curve) 적용:\n  - 가속 구간: 빠르게 도달\n  - 감속 구간: 부드럽게 정지 (Beta, Gamma 파라미터 최적화)\n• 물리적 충격량 최소화 입증",
              notes="[발표 멘트] 첫 번째 해결책은 모션 프로파일 자체를 부드럽게 만드는 비대칭 S-Curve입니다. 브레이크를 천천히 밟아 덜컹거림을 줄이는 원리입니다. [이 슬라이드에 웹 UI의 '가속도 및 저크' 차트 캡처본 삽입]")
              
    # Slide 4: 핵심 솔루션 2
    add_slide(prs, "3. 솔루션 2: Input Shaping 제어", 
              "• 파동의 상쇄 간섭 원리 활용\n• 첫 번째 충격으로 발생한 진동을, 두 번째 충격(임펄스)으로 정확히 상쇄(Cancel out)\n• 로봇 팔의 고유진동수(fn)와 감쇠비(zeta)를 계산하여 ZV/ZVD 쉐이퍼 설계",
              notes="[발표 멘트] 두 번째는 인풋 쉐이핑입니다. 노이즈 캔슬링 이어폰처럼 진동을 진동으로 덮어버리는 제어 공학의 꽃입니다. 수식 모델링을 통해 로봇에 맞춤 적용했습니다.")

    # Slide 5: 구현
    add_slide(prs, "4. 디지털 트윈 및 AI 자동 튜닝 개발", 
              "• 이론적 계산을 넘어 실제 3D 웹 환경(Digital Twin) 구현\n• 4가지 고해상도 실시간 분석 차트 탑재\n• AI 최적화 (Particle Swarm Optimization):\n  - 사람이 찾기 힘든 최적의 변수를 백그라운드 시뮬레이션으로 자동 탐색",
              notes="[발표 멘트] 단순히 수식으로 풀고 끝난 게 아닙니다. 실제 로봇 동작을 3D 시뮬레이터로 완벽 구현했으며, AI 버튼 하나만 누르면 최적의 파라미터를 스스로 찾아내도록 시스템을 구축했습니다. [이 슬라이드에 3D 풀스크린 화면 캡처 삽입]")
              
    # Slide 6: 결론
    add_slide(prs, "5. 결과 분석 및 기대효과", 
              "• [결과] 제어 미적용 시 대비 정착 시간 대폭 감소\n• [효과] 오프라인 환경에서 모터 파라미터 사전 튜닝 가능\n• [결론] 장비 셋업 시간 단축 및 실질적 UPH 향상에 기여",
              notes="[발표 멘트] 결과적으로 진동 억제를 통해 정착 시간을 획기적으로 줄였습니다. 우측 하단의 진동 차트를 보시면 미적용 대비 완전히 진동이 잡힌 것을 확인하실 수 있습니다. [이 슬라이드에 잔류 진동 차트 캡처 삽입]")

    prs.save('[4학년] (작성가이드)2026년 산학프로젝트 최종보고 발표자료_초안.pptx')

if __name__ == '__main__':
    create_poster()
    create_presentation()
    print("PPT generated.")
